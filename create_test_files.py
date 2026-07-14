import os
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas

os.makedirs("test_files", exist_ok=True)

# TXT
with open("test_files/test.txt", "w") as f:
    f.write("This is a test text file for RAG ingestion.")

# DOCX
doc = Document()
doc.add_paragraph("This is a test docx file for RAG ingestion.")
doc.save("test_files/test.docx")

# PPTX
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Test Title"
slide.shapes.placeholders[1].text = "This is a test pptx file for RAG ingestion."
prs.save("test_files/test.pptx")

# PDF
c = canvas.Canvas("test_files/test.pdf")
c.drawString(100, 750, "This is a test pdf file for RAG ingestion.")
c.save()
print("Files created!")
