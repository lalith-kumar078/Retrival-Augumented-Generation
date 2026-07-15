"""Document ingestion service — handles upload, parsing, chunking, and embedding."""

import asyncio
import hashlib
import logging
import time
import uuid
from datetime import datetime
from pathlib import Path
from typing import Optional

from app.config import settings
from app.services.chunking import chunk_document_pages
from app.services.embedding import embed_texts
from app.services.vectorstore import (
    store_chunks,
    store_document_metadata,
    update_document_status,
    document_hash_exists,
)

logger = logging.getLogger(__name__)

INGESTION_TIMEOUT_SECONDS = 300  # 5 minutes

import threading
_converter_local = threading.local()


def compute_file_hash(file_bytes: bytes) -> str:
    """Compute SHA-256 hash of file content."""
    return hashlib.sha256(file_bytes).hexdigest()


def parse_document(file_path: str, filename: str) -> list[dict]:
    """
    Parse a document, preserving structure.
    Returns a list of page dicts with 'text' and 'page_number'.

    Routing:
    - .txt, .ppt  → direct fallback (Docling doesn't support these)
    - .pdf         → direct fallback via PyMuPDF (Docling's OCR pipeline is
                     extremely slow on CPU and causes timeouts)
    - .docx, .pptx → try Docling first, fall back on error
    """
    ext = Path(filename).suffix.lower()

    # Formats that should SKIP Docling entirely
    if ext in (".txt", ".ppt", ".pdf"):
        logger.info(f"[INGEST] Using fallback parser for '{filename}' (ext={ext})")
        return _fallback_parse(file_path, filename)

    # For .docx and .pptx, try Docling first
    try:
        from docling.document_converter import DocumentConverter, PdfFormatOption
        from docling.datamodel.pipeline_options import PdfPipelineOptions, EasyOcrOptions
        
        if not hasattr(_converter_local, "converter"):
            pipeline_options = PdfPipelineOptions()
            pipeline_options.do_ocr = False
            pipeline_options.do_table_structure = False
            
            _converter_local.converter = DocumentConverter(
                format_options={"pdf": PdfFormatOption(pipeline_options=pipeline_options)}
            )
        
        logger.info("[INGEST] Docling parse started")
        result = _converter_local.converter.convert(file_path)
        
        # Extract text preserving structure
        doc = result.document
        full_text = doc.export_to_markdown()
        
        # Try to split by pages if possible
        pages = []
        if hasattr(doc, 'pages') and doc.pages:
            page_numbers = doc.pages.keys() if isinstance(doc.pages, dict) else range(1, len(doc.pages) + 1)
            for page_num in page_numbers:
                page_text = doc.export_to_markdown(page_no=page_num)
                pages.append({
                    "text": page_text,
                    "page_number": page_num
                })
        else:
            # Fallback: treat entire document as one page
            pages = [{"text": full_text, "page_number": 1}]
        
        logger.info(f"Parsed {filename} with Docling: {len(pages)} pages")
        
        full_text_extracted = "\n".join(p["text"] for p in pages)
        logger.info(f"[INGEST] Docling parse completed. Extracted text length: {len(full_text_extracted)}. Preview: {full_text_extracted[:200]}")
        return pages
        
    except ImportError:
        logger.warning("Docling not installed, using fallback parser")
        return _fallback_parse(file_path, filename)
    except Exception as e:
        logger.error(f"Docling parsing failed for {filename}: {e}")
        return _fallback_parse(file_path, filename)



def _fallback_parse(file_path: str, filename: str) -> list[dict]:
    """Fallback parser for when Docling is not available."""
    ext = Path(filename).suffix.lower()
    
    if ext == ".pdf":
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise RuntimeError("PDF parser (PyMuPDF) is not installed in the virtual environment.")
        
        try:
            doc = fitz.open(file_path)
            pages = []
            for i, page in enumerate(doc):
                # Use blocks to separate paragraphs properly with \n\n
                blocks = page.get_text("blocks")
                text_blocks = [b[4].strip() for b in blocks if b[6] == 0 and b[4].strip()]
                text = "\n\n".join(text_blocks)
                if text.strip():
                    pages.append({"text": text, "page_number": i + 1})
            doc.close()
            if not pages:
                raise ValueError("PDF document contains no extractable text.")
            return pages
        except Exception as e:
            raise RuntimeError(f"Error reading PDF file: {e}")
    
    elif ext in (".docx", ".doc"):
        try:
            import docx
        except ImportError:
            raise RuntimeError("Word document parser (python-docx) is not installed in the virtual environment.")
            
        try:
            doc = docx.Document(file_path)
            text = "\n\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            if not text.strip():
                raise ValueError("Word document contains no extractable text.")
            return [{"text": text, "page_number": 1}]
        except Exception as e:
            raise RuntimeError(f"Error reading Word document: {e}")
            
    elif ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError("PowerPoint parser (python-pptx) is not installed in the virtual environment.")
            
        try:
            prs = Presentation(file_path)
            pages = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                text = "\n\n".join(slide_text)
                if text.strip():
                    pages.append({"text": text, "page_number": i + 1})
            if not pages:
                raise ValueError("PowerPoint document contains no extractable text.")
            return pages
        except Exception as e:
            raise RuntimeError(f"Error reading PowerPoint document: {e}")
    
    elif ext == ".txt":
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            if not text.strip():
                raise ValueError("Text file is empty.")
            return [{"text": text, "page_number": 1}]
        except Exception as e:
            raise RuntimeError(f"Error reading text file: {e}")
    
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _run_pipeline_sync(doc_id: str, file_path: str, filename: str, document_type: str) -> int:
    """
    Run the CPU-bound ingestion pipeline synchronously (designed to be called via to_thread).
    Returns the number of chunks stored.
    Raises on any error — caller is responsible for updating document status.
    """
    pipeline_start = time.time()

    # ── Stage 1: Parse ──
    t0 = time.time()
    logger.info(f"[INGEST {doc_id[:8]}] Stage 1/3 PARSE — starting for '{filename}'")
    pages = parse_document(file_path, filename)
    t_parse = time.time() - t0
    logger.info(f"[INGEST {doc_id[:8]}] Stage 1/3 PARSE — done in {t_parse:.2f}s, {len(pages)} pages extracted")

    if not pages:
        raise RuntimeError("No pages extracted from document")

    # ── Stage 2: Chunk ──
    t0 = time.time()
    
    full_text_input = "\n".join(p["text"] for p in pages)
    logger.info(f"[INGEST] Chunking started. Input length: {len(full_text_input)}, type: {type(full_text_input)}")
    
    chunk_metadata = {
        "doc_id": doc_id,
        "filename": filename,
        "document_type": document_type,
        "date_uploaded": datetime.utcnow().isoformat(),
    }
    chunks = chunk_document_pages(pages, metadata=chunk_metadata)
    t_chunk = time.time() - t0
    logger.info(f"[INGEST] Chunking completed. Chunk count: {len(chunks)}")
    
    if not chunks:
        raise RuntimeError("No content extracted from document after chunking")

    # ── Stage 3a: Embed ──
    t0 = time.time()
    logger.info(f"[INGEST] Embedding started for {len(chunks)} chunks")
    texts = [c["content"] for c in chunks]
    embeddings = embed_texts(texts)
    t_embed = time.time() - t0
    logger.info("[INGEST] Embedding completed")

    # ── Stage 3b: Store in SQLite-vec ──
    t0 = time.time()
    logger.info(f"[INGEST {doc_id[:8]}] Stage 3/3 STORE — writing {len(chunks)} chunks to vector store")
    chunk_records = []
    for i, (chunk, emb) in enumerate(zip(chunks, embeddings)):
        chunk_records.append({
            "chunk_id": f"{doc_id}_chunk_{i}",
            "doc_id": doc_id,
            "filename": chunk.get("filename", filename),
            "page_number": chunk.get("page_number", 1),
            "content": chunk["content"],
            "document_type": chunk.get("document_type", document_type),
            "date_uploaded": chunk.get("date_uploaded", datetime.utcnow().isoformat()),
            "vector_f32": emb["vector_f32"],
            "vec_min": emb["vec_min"],
            "vec_max": emb["vec_max"],
        })

    logger.info("[INGEST] Writing to sqlite-vec (vector table)")
    logger.info("[INGEST] Writing to FTS5 table")
    store_chunks(chunk_records)
    t_store = time.time() - t0
    logger.info(f"[INGEST {doc_id[:8]}] Stage 3/3 STORE — done in {t_store:.2f}s")

    total_time = time.time() - pipeline_start
    logger.info(
        f"[INGEST {doc_id[:8]}] Pipeline complete for '{filename}': "
        f"{len(chunk_records)} chunks in {total_time:.2f}s "
        f"(parse={t_parse:.2f}s, chunk={t_chunk:.2f}s, embed={t_embed:.2f}s, store={t_store:.2f}s)"
    )
    return len(chunk_records)


async def run_ingestion_background(doc_id: str, file_path: str, filename: str, document_type: str):
    """
    Run the full ingestion pipeline in a background thread with a timeout.
    Updates document status to 'ready' on success or 'failed' on error/timeout.
    This is designed to be launched as a fire-and-forget background task.
    """
    try:
        logger.info(f"[INGEST {doc_id[:8]}] Background pipeline starting for '{filename}'")

        total_chunks = await asyncio.wait_for(
            asyncio.to_thread(
                _run_pipeline_sync, doc_id, file_path, filename, document_type
            ),
            timeout=INGESTION_TIMEOUT_SECONDS,
        )

        update_document_status(doc_id, "ready", total_chunks)
        logger.info("[INGEST] Document status updated to: ready")
        logger.info(f"[INGEST {doc_id[:8]}] ✓ Document '{filename}' is now READY ({total_chunks} chunks)")

    except asyncio.TimeoutError:
        logger.error(
            f"[INGEST {doc_id[:8]}] ✗ Pipeline TIMED OUT after {INGESTION_TIMEOUT_SECONDS}s for '{filename}'"
        )
        update_document_status(doc_id, "failed")

    except Exception as e:
        logger.error(f"[INGEST {doc_id[:8]}] ✗ Pipeline FAILED for '{filename}': {e}", exc_info=True)
        update_document_status(doc_id, f"failed: {str(e)}")
        logger.info(f"[INGEST] Document status updated to: failed")


async def ingest_document(
    file_bytes: bytes,
    filename: str,
    content_type: str = ""
) -> dict:
    """
    Full ingestion pipeline:
    1. Hash check (skip if duplicate)
    2. Save file to disk
    3. Store initial metadata (status=processing)
    4. Launch background pipeline (parse → chunk → embed → store)
    
    Returns document metadata dict immediately with status 'processing'.
    The background task will update status to 'ready' or 'failed'.
    """
    # Step 1: Hash check
    logger.info(f"[INGEST] File received: {filename}, size: {len(file_bytes)} bytes")
    file_hash = compute_file_hash(file_bytes)
    existing_doc_id = document_hash_exists(file_hash)
    logger.info(f"[INGEST] Hash computed: {file_hash}, duplicate check: {bool(existing_doc_id)}")
    if existing_doc_id:
        logger.info(f"Document already exists with hash {file_hash[:8]}..., doc_id={existing_doc_id}")
        return {
            "doc_id": existing_doc_id,
            "filename": filename,
            "status": "duplicate",
            "message": "Document already uploaded (hash match)",
            "total_chunks": 0
        }
    
    # Step 2: Generate doc_id and save to disk
    doc_id = str(uuid.uuid4())
    ext = Path(filename).suffix.lower()
    document_type = ext.lstrip(".")
    
    upload_dir = Path(settings.upload_abs_path)
    file_path = upload_dir / f"{doc_id}{ext}"
    file_path.write_bytes(file_bytes)
    logger.info(f"[INGEST] File saved to disk at: {file_path}")
    
    # Step 3: Store initial metadata with status=processing
    doc_meta = {
        "doc_id": doc_id,
        "filename": filename,
        "file_hash": file_hash,
        "document_type": document_type,
        "author": "",
        "date_uploaded": datetime.utcnow().isoformat(),
        "total_pages": 0,
        "total_chunks": 0,
        "status": "processing"
    }
    store_document_metadata(doc_meta)
    
    # Step 4: Launch the pipeline as a background asyncio task
    # This runs in a separate thread (via to_thread) and does NOT block the event loop.
    asyncio.create_task(
        run_ingestion_background(doc_id, str(file_path), filename, document_type)
    )
    
    logger.info(f"[INGEST {doc_id[:8]}] Upload accepted, background pipeline launched for '{filename}'")
    
    return {
        "doc_id": doc_id,
        "filename": filename,
        "status": "processing",
        "message": "Document uploaded — processing in background",
        "total_chunks": 0
    }
