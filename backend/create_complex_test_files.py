import os
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF

os.makedirs("test_files", exist_ok=True)

# 1. TXT
txt_content = """RAG Ingestion Architecture Document

The modern Retrieval-Augmented Generation (RAG) ingestion pipeline consists of several key stages:
1. Parsing: Documents (PDFs, Word docs, PowerPoint presentations) are parsed using tools like Docling or PyMuPDF to extract text and structure.
2. Chunking: The extracted text is divided into semantically meaningful chunks, typically using overlapping sliding windows to preserve context across boundaries.
3. Embedding: A specialized model, such as all-MiniLM-L6-v2, converts the chunks into dense vector representations.
4. Storage: The vectors are stored alongside metadata in a vector database like SQLite-vec, enabling efficient similarity search.

This architecture ensures that the LLM has access to highly relevant context when answering user queries.
"""
with open("test_files/test.txt", "w") as f:
    f.write(txt_content)

# 2. DOCX
doc = Document()
doc.add_heading("Employee Benefits Handbook", 0)
doc.add_heading("Health Insurance", level=1)
doc.add_paragraph("All full-time employees are eligible for the company's comprehensive health insurance plan. Coverage includes medical, dental, and vision benefits. The plan is administered by HealthCorp and covers 80% of out-of-network costs.")
doc.add_heading("Paid Time Off", level=1)
doc.add_paragraph("Employees accrue 15 days of Paid Time Off (PTO) per year during their first three years of service. After three years, this increases to 20 days per year. Unused PTO can be rolled over to the next year, up to a maximum of 10 days.")
doc.save("test_files/test.docx")

# 3. PPTX
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Q3 Financial Results Overview"
subtitle.text = "A detailed breakdown of our financial performance in the third quarter."

bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Key Metrics"
tf = body_shape.text_frame
tf.text = "Revenue grew by 25% year-over-year, reaching $50 million."
p = tf.add_paragraph()
p.text = "Operating margins expanded to 18%, driven by efficiency gains."
p = tf.add_paragraph()
p.text = "Customer acquisition costs decreased by 15% due to optimized marketing spend."
prs.save("test_files/test.pptx")

# 4. PDF (we'll generate it with ReportLab, let's install reportlab first, or just use fitz to create an empty PDF and insert text, or just create it with reportlab)
